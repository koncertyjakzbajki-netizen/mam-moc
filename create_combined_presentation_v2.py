#!/usr/bin/env python3
"""
Ulepszona wersja prezentacji łączącej oba programy.
Zmiany:
- Zajęcia pokazowe: tylko JEDNE (nie dwa)
- Dodano slajd z zajęciami pokazowymi dla Muzycznej Bajolandii
- Usunięto informacje o cenie zajęć pokazowych
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

# Kolory
PINK = RGBColor(255, 107, 107)
TEAL = RGBColor(78, 205, 196)
BLUE = RGBColor(69, 183, 209)
DARK = RGBColor(45, 52, 54)
WHITE = RGBColor(255, 255, 255)

def add_text_box(slide, left, top, width, height, text, font_size=18, bold=False, color=DARK, align=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.alignment = align
    return txBox

def create_presentation():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # ========== SLIDE 1: Tytuł ==========
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    background.fill.solid()
    background.fill.fore_color.rgb = RGBColor(45, 52, 54)
    background.line.fill.background()
    
    add_text_box(slide, Inches(0.5), Inches(1.5), Inches(12.3), Inches(1.2),
                 "Dwa wyjątkowe programy dla Waszych dzieci", 
                 font_size=44, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text_box(slide, Inches(0.5), Inches(2.8), Inches(12.3), Inches(0.8),
                 "Rok szkolny 2026/2027 • Wojszycka Akademia Talentów", 
                 font_size=24, color=RGBColor(200, 200, 200), align=PP_ALIGN.CENTER)
    
    box1 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1), Inches(4), Inches(5), Inches(2.5))
    box1.fill.solid()
    box1.fill.fore_color.rgb = PINK
    box1.line.fill.background()
    add_text_box(slide, Inches(1.2), Inches(4.2), Inches(4.6), Inches(0.6),
                 "Mam Moc!", font_size=28, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text_box(slide, Inches(1.2), Inches(4.8), Inches(4.6), Inches(1.5),
                 "Superbohaterowie Małych Sukcesów\nRozwój pewności siebie, wytrwałości\ni inteligencji emocjonalnej", 
                 font_size=16, color=WHITE, align=PP_ALIGN.CENTER)
    
    box2 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(7.3), Inches(4), Inches(5), Inches(2.5))
    box2.fill.solid()
    box2.fill.fore_color.rgb = TEAL
    box2.line.fill.background()
    add_text_box(slide, Inches(7.5), Inches(4.2), Inches(4.6), Inches(0.6),
                 "Muzyczna Bajolandia", font_size=28, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text_box(slide, Inches(7.5), Inches(4.8), Inches(4.6), Inches(1.5),
                 "Śpiewaj, graj i baw się razem z nami\nRozwój przez muzykę, rytm i ruch", 
                 font_size=16, color=WHITE, align=PP_ALIGN.CENTER)

    # ========== SLIDE 2: Porównanie ==========
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_text_box(slide, Inches(0.5), Inches(0.3), Inches(12.3), Inches(0.8),
                 "Który program wybrać? – Porównanie", 
                 font_size=32, bold=True, color=DARK, align=PP_ALIGN.CENTER)
    
    add_text_box(slide, Inches(0.5), Inches(1.3), Inches(4), Inches(0.5),
                 "Cecha", font_size=18, bold=True, color=DARK)
    add_text_box(slide, Inches(4.7), Inches(1.3), Inches(4), Inches(0.5),
                 "Mam Moc!", font_size=18, bold=True, color=PINK)
    add_text_box(slide, Inches(8.9), Inches(1.3), Inches(4), Inches(0.5),
                 "Muzyczna Bajolandia", font_size=18, bold=True, color=TEAL)
    
    rows = [
        ("Główny cel", "Rozwój charakteru i kompetencji miękkich", "Rozwój przez muzykę, śpiew i rytm"),
        ("Największe korzyści", "Wytrwałość, pewność siebie, radzenie z emocjami", "Rozwój mowy, koordynacji, kreatywności"),
        ("Metody", "Bajki z pacynkami, role-playing, ruch", "Śpiew, instrumenty, taniec, rytm"),
        ("Dla kogo szczególnie", "Dzieci potrzebujące pewności siebie i wytrwałości", "Dzieci, które lubią śpiew i muzykę"),
        ("Czas trwania", "30 min / tydzień", "30 min / tydzień"),
        ("Liczba dzieci", "max 10", "max 10"),
        ("Cena", "25–33 zł / zajęcia", "120 zł / miesiąc"),
    ]
    
    y_pos = 1.9
    for i, (cecha, mam_moc, bajolandia) in enumerate(rows):
        bg_color = RGBColor(245, 245, 245) if i % 2 == 0 else WHITE
        row_bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(y_pos), Inches(12.3), Inches(0.65))
        row_bg.fill.solid()
        row_bg.fill.fore_color.rgb = bg_color
        row_bg.line.fill.background()
        
        add_text_box(slide, Inches(0.6), Inches(y_pos + 0.1), Inches(3.8), Inches(0.5),
                     cecha, font_size=14, bold=True, color=DARK)
        add_text_box(slide, Inches(4.8), Inches(y_pos + 0.1), Inches(3.8), Inches(0.5),
                     mam_moc, font_size=13, color=DARK)
        add_text_box(slide, Inches(9), Inches(y_pos + 0.1), Inches(3.8), Inches(0.5),
                     bajolandia, font_size=13, color=DARK)
        y_pos += 0.65

    # ========== SLIDE 3: Mam Moc - Szczegóły ==========
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.2))
    header.fill.solid()
    header.fill.fore_color.rgb = PINK
    header.line.fill.background()
    
    add_text_box(slide, Inches(0.5), Inches(0.3), Inches(12.3), Inches(0.7),
                 "Mam Moc! – Superbohaterowie Małych Sukcesów", 
                 font_size=28, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text_box(slide, Inches(0.5), Inches(1.5), Inches(12.3), Inches(0.6),
                 "Program rozwijający pewność siebie, wytrwałość i inteligencję emocjonalną", 
                 font_size=18, color=DARK, align=PP_ALIGN.CENTER)
    
    benefits = [
        "• Bardziej wytrwałe – nie poddaje się przy trudnościach",
        "• Pewniejsze siebie i odważniejsze w próbowaniu nowych rzeczy",
        "• Lepsze w radzeniu sobie z emocjami i frustracją",
        "• Szczęśliwsze i dumne ze siebie",
        "• Silniejsze wewnętrznie i gotowe na wyzwania szkoły i życia"
    ]
    y = 2.3
    for benefit in benefits:
        add_text_box(slide, Inches(0.8), Inches(y), Inches(11.5), Inches(0.5), benefit, font_size=18, color=DARK)
        y += 0.55
    
    add_text_box(slide, Inches(0.5), Inches(5.5), Inches(12.3), Inches(1.2),
                 "Metody: Autorskie bajki z pacynkami (Jeżyk, Sowa, Żabka) • Wcielanie się w role\nUnikalne piosenki • Ruch i zabawa • Mała grupa (max 10 dzieci)", 
                 font_size=16, color=RGBColor(100, 100, 100), align=PP_ALIGN.CENTER)

    # ========== SLIDE 4: Muzyczna Bajolandia - Szczegóły ==========
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.2))
    header.fill.solid()
    header.fill.fore_color.rgb = TEAL
    header.line.fill.background()
    
    add_text_box(slide, Inches(0.5), Inches(0.3), Inches(12.3), Inches(0.7),
                 "Muzyczna Bajolandia – Śpiewaj, graj i baw się razem z nami", 
                 font_size=26, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text_box(slide, Inches(0.5), Inches(1.5), Inches(12.3), Inches(0.6),
                 "Program rozwijający przez muzykę, rytm, śpiew i ruch", 
                 font_size=18, color=DARK, align=PP_ALIGN.CENTER)
    
    benefits2 = [
        "• Rozwój mowy i słownictwa poprzez śpiew",
        "• Poprawa koordynacji, rytmu i świadomości ciała",
        "• Budowanie pewności siebie przez występy i wspólne muzykowanie",
        "• Rozwój kreatywności i wyobraźni",
        "• Wzmacnianie koncentracji i pamięci (trening dla mózgu)"
    ]
    y = 2.3
    for benefit in benefits2:
        add_text_box(slide, Inches(0.8), Inches(y), Inches(11.5), Inches(0.5), benefit, font_size=18, color=DARK)
        y += 0.55
    
    add_text_box(slide, Inches(0.5), Inches(5.5), Inches(12.3), Inches(1.2),
                 "Metody: Śpiew • Gra na instrumentach (cymbałki, trójkąty, grzechotki) • Taniec i ruch\nGry rytmiczne • Mała grupa (max 10 dzieci) • Ciepła, radosna atmosfera", 
                 font_size=16, color=RGBColor(100, 100, 100), align=PP_ALIGN.CENTER)

    # ========== SLIDE 5: Który wybrać? ==========
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_text_box(slide, Inches(0.5), Inches(0.3), Inches(12.3), Inches(0.8),
                 "Który program wybrać?", 
                 font_size=32, bold=True, color=DARK, align=PP_ALIGN.CENTER)
    
    box_mam = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(1.3), Inches(6), Inches(5.5))
    box_mam.fill.solid()
    box_mam.fill.fore_color.rgb = RGBColor(255, 240, 240)
    box_mam.line.color.rgb = PINK
    
    add_text_box(slide, Inches(0.7), Inches(1.5), Inches(5.6), Inches(0.6),
                 "Wybierz „Mam Moc!", font_size=22, bold=True, color=PINK, align=PP_ALIGN.CENTER)
    add_text_box(slide, Inches(0.7), Inches(2.2), Inches(5.6), Inches(4.3),
                 "Jeśli zależy Wam szczególnie na:\n\n• Rozwoju pewności siebie\n• Wytrwałości i radzeniu sobie z trudnościami\n• Inteligencji emocjonalnej\n• Umiejętności współpracy\n\nTo program, który realnie zmienia sposób, w jaki dziecko podchodzi do wyzwań.", 
                 font_size=15, color=DARK)
    
    box_bajo = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.8), Inches(1.3), Inches(6), Inches(5.5))
    box_bajo.fill.solid()
    box_bajo.fill.fore_color.rgb = RGBColor(230, 250, 245)
    box_bajo.line.color.rgb = TEAL
    
    add_text_box(slide, Inches(7), Inches(1.5), Inches(5.6), Inches(0.6),
                 "Wybierz „Muzyczną Bajolandię", font_size=22, bold=True, color=TEAL, align=PP_ALIGN.CENTER)
    add_text_box(slide, Inches(7), Inches(2.2), Inches(5.6), Inches(4.3),
                 "Jeśli zależy Wam szczególnie na:\n\n• Rozwoju mowy i słownictwa\n• Koordynacji i rytmie\n• Kreatywności i wyobraźni\n• Radości z muzyki i wspólnego śpiewania\n\nTo program, który rozwija przez zabawę muzyczną.", 
                 font_size=15, color=DARK)

    # ========== SLIDE 6: Zajęcia pokazowe - Mam Moc ==========
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.2))
    header.fill.solid()
    header.fill.fore_color.rgb = RGBColor(69, 183, 209)
    header.line.fill.background()
    
    add_text_box(slide, Inches(0.5), Inches(0.3), Inches(12.3), Inches(0.7),
                 "Zajęcia pokazowe – „Mam Moc!", 
                 font_size=28, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    
    add_text_box(slide, Inches(0.5), Inches(1.5), Inches(12.3), Inches(1.5),
                 "Chcesz zobaczyć, jak Twoje dziecko może rozkwitnąć?\nZapisz je na zajęcia pokazowe i zobacz na własne oczy!", 
                 font_size=20, color=DARK, align=PP_ALIGN.CENTER)
    
    add_text_box(slide, Inches(0.8), Inches(3.2), Inches(11.5), Inches(3.5),
                 "Na zajęciach pokazowych dziecko będzie miało okazję:\n\n• Przeżyć małe sukcesy i poczuć dumę\n• Zobaczyć, jak wiele radości może dawać nauka przez zabawę\n• Poczuć, że jest ważne i doceniane\n\nPo zajęciach spokojnie zdecydujesz, czy chcesz kontynuować cały rok.", 
                 font_size=17, color=DARK)

    # ========== SLIDE 7: Zajęcia pokazowe - Muzyczna Bajolandia (NOWY) ==========
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.2))
    header.fill.solid()
    header.fill.fore_color.rgb = TEAL
    header.line.fill.background()
    
    add_text_box(slide, Inches(0.5), Inches(0.3), Inches(12.3), Inches(0.7),
                 "Zajęcia pokazowe – „Muzyczna Bajolandia", 
                 font_size=28, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    
    add_text_box(slide, Inches(0.5), Inches(1.5), Inches(12.3), Inches(1.5),
                 "Chcesz zobaczyć, jak muzyka może rozwijać Wasze dziecko?\nZapisz je na zajęcia pokazowe i przekonaj się na własne oczy!", 
                 font_size=20, color=DARK, align=PP_ALIGN.CENTER)
    
    add_text_box(slide, Inches(0.8), Inches(3.2), Inches(11.5), Inches(3.5),
                 "Na zajęciach pokazowych dziecko będzie miało okazję:\n\n• Śpiewać, grać na instrumentach i tańczyć w radosnej atmosferze\n• Doświadczyć, jak muzyka pomaga w rozwoju mowy i koordynacji\n• Poczuć radość ze wspólnego muzykowania z rówieśnikami\n\nPo zajęciach spokojnie zdecydujesz, czy chcecie kontynuować cały rok.", 
                 font_size=17, color=DARK)

    # ========== SLIDE 8: Ceny ==========
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_text_box(slide, Inches(0.5), Inches(0.3), Inches(12.3), Inches(0.7),
                 "Ceny i warunki", 
                 font_size=32, bold=True, color=DARK, align=PP_ALIGN.CENTER)
    
    box1 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(1.2), Inches(6), Inches(5.8))
    box1.fill.solid()
    box1.fill.fore_color.rgb = RGBColor(255, 245, 245)
    box1.line.color.rgb = PINK
    
    add_text_box(slide, Inches(0.7), Inches(1.4), Inches(5.6), Inches(0.5),
                 "Mam Moc!", font_size=24, bold=True, color=PINK, align=PP_ALIGN.CENTER)
    add_text_box(slide, Inches(0.7), Inches(2), Inches(5.6), Inches(4.8),
                 "33 zł – płatność pojedyncza\n30 zł – płatność miesięczna\n25 zł – przy płatności z góry / w ratach\n\nPromocja przy zapisie do końca lipca/sierpnia:\n\nI rata (do 31.08) → 450 zł\nII rata (do 31.12) → 450 zł\n\nRazem tylko 900 zł za cały rok\n(zamiast 1188 zł)", 
                 font_size=16, color=DARK)
    
    box2 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.8), Inches(1.2), Inches(6), Inches(5.8))
    box2.fill.solid()
    box2.fill.fore_color.rgb = RGBColor(235, 250, 245)
    box2.line.color.rgb = TEAL
    
    add_text_box(slide, Inches(7), Inches(1.4), Inches(5.6), Inches(0.5),
                 "Muzyczna Bajolandia", font_size=24, bold=True, color=TEAL, align=PP_ALIGN.CENTER)
    add_text_box(slide, Inches(7), Inches(2), Inches(5.6), Inches(4.8),
                 "120 zł / miesiąc\n\nW cenie:\n• 4 zajęcia miesięczne (30 min)\n• Mała grupa (max 10 dzieci)\n• Wszystkie instrumenty na miejscu\n• Ciepła, indywidualna atmosfera\n\nTo jedna z najtańszych inwestycji\nw rozwój emocjonalny, społeczny\ni intelektualny dziecka.", 
                 font_size=16, color=DARK)

    # ========== SLIDE 9: Kontakt / Ankieta ==========
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    background.fill.solid()
    background.fill.fore_color.rgb = RGBColor(45, 52, 54)
    background.line.fill.background()
    
    add_text_box(slide, Inches(0.5), Inches(1), Inches(12.3), Inches(1),
                 "Które zajęcia najbardziej Was interesują?", 
                 font_size=32, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    
    add_text_box(slide, Inches(0.5), Inches(2.2), Inches(12.3), Inches(3.5),
                 "A) „Mam Moc! – Superbohaterowie Małych Sukcesów"\n(Rozwój pewności siebie, wytrwałości i inteligencji emocjonalnej)\n\nB) „Muzyczna Bajolandia"\n(Śpiew, gra na instrumentach, taniec i rytm)\n\nC) Trudno mi powiedzieć / chciał(a)bym zobaczyć obie opcje", 
                 font_size=20, color=WHITE, align=PP_ALIGN.CENTER)
    
    add_text_box(slide, Inches(0.5), Inches(6), Inches(12.3), Inches(1),
                 "Ilona Kalińska  •  514 487 616  •  koncertyjakzbajki@gmail.com", 
                 font_size=18, color=RGBColor(200, 200, 200), align=PP_ALIGN.CENTER)

    prs.save('Polaczone_Programy_Prezentacja_v2.pptx')
    print("Ulepszona prezentacja utworzona: Polaczone_Programy_Prezentacja_v2.pptx")

if __name__ == "__main__":
    create_presentation()
